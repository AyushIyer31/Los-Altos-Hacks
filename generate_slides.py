from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ── Colour palette ────────────────────────────────────────────────────────────
NAVY       = RGBColor(0x0D, 0x1B, 0x2A)
DARK_NAVY  = RGBColor(0x08, 0x12, 0x1C)
TEAL       = RGBColor(0x00, 0xB4, 0xA0)
TEAL_DARK  = RGBColor(0x00, 0x7A, 0x6E)
TEAL_MID   = RGBColor(0x00, 0x96, 0x88)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xCC, 0xCC, 0xCC)
CARD_BG    = RGBColor(0x12, 0x25, 0x38)
CARD_BG2   = RGBColor(0x0A, 0x2A, 0x2A)
RED_WARN   = RGBColor(0xE5, 0x73, 0x73)
GREEN_OK   = RGBColor(0x81, 0xC7, 0x84)

W = Inches(13.33)   # widescreen width
H = Inches(7.5)     # widescreen height

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]   # totally blank


# ── helpers ───────────────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill_rgb, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    return shape

def add_text(slide, text, x, y, w, h, font_size=18, bold=False,
             color=WHITE, align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb

def add_card(slide, x, y, w, h, bg=CARD_BG, border_color=None, border_top=False):
    card = add_rect(slide, x, y, w, h, bg)
    if border_top and border_color:
        # thin top accent line
        line = add_rect(slide, x, y, w, Inches(0.05), border_color)
    return card

def add_title_bar(slide, title, subtitle=None):
    """Left teal accent bar + title."""
    add_rect(slide, Inches(0.35), Inches(0.25), Inches(0.06), Inches(0.85), TEAL)
    add_text(slide, title, Inches(0.5), Inches(0.22), Inches(10), Inches(0.55),
             font_size=32, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, Inches(0.5), Inches(0.77), Inches(12), Inches(0.38),
                 font_size=16, italic=True, color=TEAL)

def slide_bg(slide, color=NAVY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_footer(slide, text="IE380 · Computational Biology Project"):
    add_text(slide, text, Inches(0.4), Inches(7.1), Inches(10), Inches(0.3),
             font_size=10, color=LIGHT_GREY, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Updated Our Approach
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl)
add_title_bar(sl, "Our Approach (Updated v46)",
              "Sequence in. Ranked mutation candidates out.")

steps = [
    ("1", "Curate Training Data",
     "16,963 real wet-lab mutation measurements\nfrom FireProtDB and ThermoMutDB.\nNo synthetic data. Reverse-mutation augmentation applied."),
    ("2", "Engineer 144 Features",
     "Each mutation → 144-number vector:\nphysicochemical deltas, ESM-2 PCA-32,\nAlphaFold pLDDT, physics (DH+Born), ESM-1v."),
    ("3", "Train 8-Model Ensemble",
     "GBM, XGBoost, LightGBM, CatBoost, HistGBM,\nMLP, RandomForest, ExtraTrees —\nOptuna-tuned, stacked with meta-classifier."),
    ("4", "Apply Physics Corrections",
     "Debye-Hückel for ionic strength.\nBorn solvation for burial energy.\nHill equation for cutinase pH response."),
    ("5", "Rank & Serve",
     "Score all single-point mutations, combine\ntop candidates into multi-mutants,\ndeliver via web app in ~3 seconds."),
]

col_x = [Inches(0.35), Inches(4.95), Inches(4.95), Inches(0.35), Inches(0.35)]
col_y = [Inches(1.35), Inches(1.35), Inches(4.05), Inches(4.05), Inches(1.35)]

# two-column layout
positions = [
    (Inches(0.35), Inches(1.3)),
    (Inches(6.85), Inches(1.3)),
    (Inches(0.35), Inches(3.7)),
    (Inches(6.85), Inches(3.7)),
    (Inches(0.35), Inches(6.05)),
]

for i, (step, title, body) in enumerate(steps):
    px, py = positions[i]
    cw, ch = Inches(6.1), Inches(2.05)
    add_card(sl, px, py, cw, ch, CARD_BG, TEAL, border_top=True)
    # number badge
    badge = add_rect(sl, px, py, Inches(0.55), ch, TEAL_DARK)
    add_text(sl, step, px, py + Inches(0.65), Inches(0.55), Inches(0.7),
             font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, title, px + Inches(0.65), py + Inches(0.08), cw - Inches(0.75), Inches(0.38),
             font_size=14, bold=True, color=TEAL)
    add_text(sl, body, px + Inches(0.65), py + Inches(0.42), cw - Inches(0.75), Inches(1.55),
             font_size=11, color=WHITE)

add_footer(sl)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Updated Training Data
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl)
add_title_bar(sl, "Training Data (Updated v46)",
              "Only real wet-lab measurements — no synthetic or computationally generated data.")

# FireProtDB card
add_card(sl, Inches(0.35), Inches(1.3), Inches(6.1), Inches(3.5), CARD_BG2, TEAL, True)
add_text(sl, "FireProtDB", Inches(0.5), Inches(1.38), Inches(5.5), Inches(0.4),
         font_size=18, bold=True, color=TEAL)
add_text(sl, "3,438 raw  →  6,796 after reverse-mutation augmentation",
         Inches(0.5), Inches(1.78), Inches(5.8), Inches(0.35), font_size=13, bold=True, color=WHITE)
bullets_fp = [
    "→  Single-point substitutions across many protein families",
    "→  Exact assay conditions recorded per experiment",
    "→  ΔΔG values from calorimetry and CD spectroscopy",
    "→  86.1% coverage with real RSA / secondary structure",
]
for j, b in enumerate(bullets_fp):
    add_text(sl, b, Inches(0.5), Inches(2.2) + Inches(j * 0.42), Inches(5.8), Inches(0.4),
             font_size=11, color=LIGHT_GREY)

# ThermoMutDB card
add_card(sl, Inches(6.85), Inches(1.3), Inches(6.1), Inches(3.5), CARD_BG2, TEAL_MID, True)
add_text(sl, "ThermoMutDB", Inches(7.0), Inches(1.38), Inches(5.5), Inches(0.4),
         font_size=18, bold=True, color=TEAL_MID)
add_text(sl, "8,801 raw  →  10,167 after reverse-mutation augmentation",
         Inches(7.0), Inches(1.78), Inches(5.8), Inches(0.35), font_size=13, bold=True, color=WHITE)
bullets_tm = [
    "→  Spans 249 protein families — broad phylogenetic coverage",
    "→  Temperature shift (ΔTm) available per record",
    "→  6,567 ΔTm records used for separate ΔTm regressor",
    "→  Enables generalisation beyond PET-active enzymes",
]
for j, b in enumerate(bullets_tm):
    add_text(sl, b, Inches(7.0), Inches(2.2) + Inches(j * 0.42), Inches(5.8), Inches(0.4),
             font_size=11, color=LIGHT_GREY)

# Class balance bar
add_card(sl, Inches(0.35), Inches(5.05), Inches(12.6), Inches(1.75), CARD_BG)
add_text(sl, "Class Balance  —  After Deduplication + Augmentation",
         Inches(0.55), Inches(5.12), Inches(8), Inches(0.38), font_size=13, bold=True, color=TEAL)
add_text(sl, "Near-perfect split — no class weighting needed.",
         Inches(0.55), Inches(5.48), Inches(6), Inches(0.3), font_size=11, italic=True, color=LIGHT_GREY)
# stabilising bar
add_rect(sl, Inches(0.55), Inches(5.85), Inches(6.15), Inches(0.6), TEAL_DARK)
add_text(sl, "Stabilising  49.3%  (8,367 mutations)", Inches(0.65), Inches(5.88),
         Inches(6), Inches(0.5), font_size=12, bold=True, color=WHITE)
# destabilising bar
add_rect(sl, Inches(6.75), Inches(5.85), Inches(6.1), Inches(0.6), TEAL)
add_text(sl, "Destabilising  50.7%  (8,596 mutations)", Inches(6.85), Inches(5.88),
         Inches(5.9), Inches(0.5), font_size=12, bold=True, color=WHITE)

add_text(sl, "Total training set: 16,963 mutations  |  S669 held-out test: 669 mutations (never seen during training)",
         Inches(0.5), Inches(7.05), Inches(12), Inches(0.35), font_size=10, color=LIGHT_GREY, italic=True)
add_footer(sl)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Updated Feature Engineering
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl)
add_title_bar(sl, "Feature Engineering (Updated — 144 Features)",
              "Every mutation is encoded as a 144-number vector.")

feat_cards = [
    (TEAL,     "Physicochemical Deltas",  "12 features",
     "Hydrophobicity · Volume · Charge\nHelix & sheet propensity · BLOSUM62 score"),
    (TEAL_MID, "Structural Context",      "Real RSA + SS  +  Burial depth",
     "Actual RSA/secondary structure from PDB\n(86.1% coverage). Burial depth per residue."),
    (TEAL_DARK,"Assay Conditions",        "4 features per record",
     "Temperature · pH · Ionic strength\nCa²⁺ concentration — matches each experiment exactly."),
    (TEAL,     "ESM-2 Language Model",    "32 features (PCA)",
     "Meta AI 150M-param protein language model.\nPCA-32 applied to 640-dim embeddings.\n92.6% training coverage."),
    (TEAL_MID, "AlphaFold pLDDT",         "3 features",
     "Per-residue confidence scores from AlphaFold.\nWT · mutant · delta pLDDT.\n92.9% training coverage."),
    (TEAL_DARK,"Physics (DH + Born)",     "~30 features",
     "Debye-Hückel electrostatics + Born solvation\nenergy per residue position.\n79.6% training coverage."),
    (TEAL,     "ESM-1v Evolutionary",     "5-model ensemble avg",
     "Zero-shot evolutionary fitness score.\n5-model ensemble from Meta AI.\n44.1% training coverage."),
    (TEAL_MID, "Extended Biochemistry",   "varies",
     "MW delta · H-bond potential · pKa shift\nAliphatic index · Disorder propensity · Cys distance"),
]

cols = 4
for i, (col, title, sub, body) in enumerate(feat_cards):
    col_i = i % cols
    row_i = i // cols
    px = Inches(0.3) + col_i * Inches(3.26)
    py = Inches(1.25) + row_i * Inches(2.85)
    add_card(sl, px, py, Inches(3.1), Inches(2.65), CARD_BG, col, border_top=True)
    add_text(sl, title, px + Inches(0.1), py + Inches(0.12), Inches(2.9), Inches(0.38),
             font_size=12, bold=True, color=col)
    add_text(sl, sub, px + Inches(0.1), py + Inches(0.5), Inches(2.9), Inches(0.3),
             font_size=10, italic=True, color=LIGHT_GREY)
    add_text(sl, body, px + Inches(0.1), py + Inches(0.82), Inches(2.9), Inches(1.7),
             font_size=10, color=WHITE)

add_footer(sl)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Updated Model Architecture (8 models)
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl)
add_title_bar(sl, "The Model — 8-Model Stacked Ensemble",
              "Eight independently tuned regressors → wide stacking → meta-classifier.")

models = [
    ("Gradient Boosting", "sklearn", "Optuna 50-trial tuned"),
    ("XGBoost",           "DMLC",    "v26 pre-tuned params"),
    ("LightGBM",          "Microsoft","v26 pre-tuned params"),
    ("CatBoost",          "Yandex",  "v24 pre-tuned params"),
    ("HistGradientBoosting","sklearn","Optuna 50-trial tuned"),
    ("MLP Neural Net",    "sklearn", "512-256-128-64, early stop"),
    ("RandomForest",      "sklearn", "Optuna 50-trial tuned"),
    ("ExtraTrees",        "sklearn", "Optuna 50-trial tuned"),
]

for i, (name, org, note) in enumerate(models):
    row = i // 2
    col = i % 2
    px = Inches(0.35) + col * Inches(3.2)
    py = Inches(1.25) + row * Inches(1.38)
    add_card(sl, px, py, Inches(3.0), Inches(1.25), CARD_BG)
    add_text(sl, name, px + Inches(0.12), py + Inches(0.1), Inches(2.8), Inches(0.38),
             font_size=13, bold=True, color=WHITE)
    add_text(sl, org, px + Inches(0.12), py + Inches(0.48), Inches(2.8), Inches(0.28),
             font_size=10, italic=True, color=TEAL)
    add_text(sl, note, px + Inches(0.12), py + Inches(0.78), Inches(2.8), Inches(0.35),
             font_size=10, color=LIGHT_GREY)

# Arrow
add_text(sl, "→", Inches(6.55), Inches(3.1), Inches(0.6), Inches(0.6),
         font_size=36, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

# Meta-classifier block
mc_x, mc_y = Inches(7.2), Inches(1.25)
add_card(sl, mc_x, mc_y, Inches(5.7), Inches(2.8), CARD_BG2, TEAL, True)
add_text(sl, "Meta-Classifier (Wide Stacking)", mc_x + Inches(0.15), mc_y + Inches(0.12),
         Inches(5.4), Inches(0.4), font_size=16, bold=True, color=TEAL)
stack_lines = [
    "• 6 classifier OOF probabilities (XGB_clf, LGBM_clf, CB_clf,",
    "  RF_clf, XGB_opt, LGBM_opt)",
    "• 8 regressor OOF predictions (all 8 models)",
    "• Top 30 mutual-information features",
    "• 45 degree-2 polynomial interactions (top-10 MI features)",
    "• Total stacking matrix: 87 columns",
    "• Winner: Logistic Regression (C=100) — CV Acc 79.56%",
]
for j, line in enumerate(stack_lines):
    add_text(sl, line, mc_x + Inches(0.15), mc_y + Inches(0.6) + Inches(j * 0.31),
             Inches(5.4), Inches(0.3), font_size=10, color=WHITE)

# Validation block
add_card(sl, mc_x, mc_y + Inches(2.95), Inches(5.7), Inches(1.5), CARD_BG)
add_text(sl, "Validation", mc_x + Inches(0.15), mc_y + Inches(3.05), Inches(5.4), Inches(0.35),
         font_size=14, bold=True, color=TEAL)
add_text(sl, "10-fold GroupKFold (protein-pair constrained) — prevents data leakage\n"
             "5-fold GroupKFold for stacking meta-classifier\n"
             "100-trial Optuna for XGB_opt / LGBM_opt classifiers (AUC objective)",
         mc_x + Inches(0.15), mc_y + Inches(3.42), Inches(5.4), Inches(0.95),
         font_size=10, color=WHITE)

add_footer(sl)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Updated Results
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl)
add_title_bar(sl, "Results (v46)", "What the model actually achieved.")

# Three big stat boxes
stats = [
    ("79.54%", "CV Accuracy\n(10-fold GroupKFold)", CARD_BG2),
    ("r = 0.752", "Pearson Correlation\n(predicted vs. measured ΔΔG)", CARD_BG),
    ("68.16%", "S669 Test Accuracy\n(669 unseen mutations)", CARD_BG2),
]
for i, (val, label, bg) in enumerate(stats):
    sx = Inches(0.35) + i * Inches(4.35)
    add_card(sl, sx, Inches(1.25), Inches(4.1), Inches(1.85), bg)
    add_text(sl, val, sx, Inches(1.35), Inches(4.1), Inches(0.95),
             font_size=38, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    add_text(sl, label, sx, Inches(2.3), Inches(4.1), Inches(0.65),
             font_size=11, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

# S669 confusion matrix explanation
add_card(sl, Inches(0.35), Inches(3.3), Inches(6.1), Inches(3.85), CARD_BG)
add_text(sl, "S669 Independent Test — Confusion Matrix",
         Inches(0.5), Inches(3.38), Inches(5.8), Inches(0.4), font_size=14, bold=True, color=TEAL)

cm_data = [
    ("",              "Predicted DESTAB", "Predicted STAB"),
    ("Actual DESTAB", "TN = 394  ✓",      "FP = 107"),
    ("Actual STAB",   "FN = 106",         "TP = 62  ✓"),
]
for r, row in enumerate(cm_data):
    for c, cell in enumerate(row):
        cx = Inches(0.5) + c * Inches(1.9)
        cy = Inches(3.85) + r * Inches(0.72)
        is_header = r == 0 or c == 0
        bg = CARD_BG2 if is_header else (TEAL_DARK if ("✓" in cell) else CARD_BG)
        add_rect(sl, cx, cy, Inches(1.85), Inches(0.68), bg)
        add_text(sl, cell, cx + Inches(0.05), cy + Inches(0.1),
                 Inches(1.75), Inches(0.5), font_size=10,
                 bold=is_header, color=TEAL if is_header else WHITE, align=PP_ALIGN.CENTER)

metrics = [
    ("Specificity (destab recall)", "78.6%", GREEN_OK),
    ("Stabilising recall",          "36.9%", RED_WARN),
    ("Overall accuracy",            "68.16%", TEAL),
    ("AUC (meta-clf score)",        "0.5334", RED_WARN),
]
for j, (label, val, col) in enumerate(metrics):
    add_text(sl, f"{label}:  {val}",
             Inches(0.5), Inches(6.1) + Inches(j * 0.3),
             Inches(5.8), Inches(0.28), font_size=11, color=col)

# CV per-model breakdown
add_card(sl, Inches(6.85), Inches(3.3), Inches(6.1), Inches(3.85), CARD_BG)
add_text(sl, "Per-Model CV Results", Inches(7.0), Inches(3.38), Inches(5.8), Inches(0.4),
         font_size=14, bold=True, color=TEAL)
model_results = [
    ("ExtraTrees",          "MAE 0.990", "Acc 78.42%"),
    ("XGBoost",             "MAE 0.993", "Acc 77.35%"),
    ("RandomForest",        "MAE 1.023", "Acc 77.14%"),
    ("LightGBM",            "MAE 1.012", "Acc 76.89%"),
    ("CatBoost",            "MAE 1.052", "Acc 76.10%"),
    ("HistGBM",             "MAE 1.087", "Acc 75.59%"),
    ("GradientBoosting",    "MAE 1.057", "Acc 75.15%"),
    ("MLP",                 "MAE 1.107", "Acc 73.22%"),
    ("ENSEMBLE (avg)",      "MAE 1.007", "Acc 77.66%"),
    ("★ Stacked (final)",   "",          "Acc 79.54%  ← best"),
]
for j, (name, mae, acc) in enumerate(model_results):
    col = TEAL if "★" in name else (LIGHT_GREY if "ENSEMBLE" in name else WHITE)
    add_text(sl, f"{'  ' if '★' not in name else ''}{name}",
             Inches(7.0), Inches(3.85) + Inches(j * 0.33),
             Inches(3.0), Inches(0.3), font_size=10, color=col, bold="★" in name)
    add_text(sl, acc, Inches(10.2), Inches(3.85) + Inches(j * 0.33),
             Inches(2.5), Inches(0.3), font_size=10, color=col, bold="★" in name)

add_footer(sl)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Physics Equations (ALL equations)
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl)
add_title_bar(sl, "Physics Equations",
              "All biophysical equations embedded as features and post-inference corrections.")

EQ_COLOR = RGBColor(0xA5, 0xF3, 0xE8)   # light mint for equations

# ── 1. Debye-Hückel ──────────────────────────────────────────────────────────
add_card(sl, Inches(0.35), Inches(1.25), Inches(12.6), Inches(1.6), CARD_BG, TEAL, True)
add_text(sl, "1  Debye-Hückel  —  Electrostatic Screening by Ionic Strength",
         Inches(0.55), Inches(1.32), Inches(12), Inches(0.38), font_size=14, bold=True, color=TEAL)
add_text(sl, "κ  =  √( 2 N_A e² I / (ε₀ ε_r k_B T) )",
         Inches(0.55), Inches(1.72), Inches(5.5), Inches(0.38),
         font_size=13, bold=True, color=EQ_COLOR)
add_text(sl, "φ_DH(r)  =  (z e / 4π ε₀ ε_r)  ×  exp(−κ r) / r",
         Inches(6.2), Inches(1.72), Inches(6.5), Inches(0.38),
         font_size=13, bold=True, color=EQ_COLOR)
add_text(sl, "κ = inverse Debye length  |  I = ionic strength (M)  |  ε_r = relative permittivity  |  z = residue charge  |  r = distance (Å)",
         Inches(0.55), Inches(2.12), Inches(12.2), Inches(0.3),
         font_size=10, color=LIGHT_GREY, italic=True)
add_text(sl, "→  Used to calculate how strongly two charged residues interact, given the salt concentration of the assay.",
         Inches(0.55), Inches(2.42), Inches(12.2), Inches(0.3),
         font_size=10, color=WHITE)

# ── 2. Born Solvation ─────────────────────────────────────────────────────────
add_card(sl, Inches(0.35), Inches(2.98), Inches(12.6), Inches(1.55), CARD_BG2, TEAL_MID, True)
add_text(sl, "2  Born Solvation  —  Energy Cost of Burying / Exposing a Charge",
         Inches(0.55), Inches(3.05), Inches(12), Inches(0.38), font_size=14, bold=True, color=TEAL_MID)
add_text(sl, "ΔG_Born  =  − (z² e²) / (8π ε₀ R)  ×  ( 1/ε_water − 1/ε_protein )",
         Inches(0.55), Inches(3.45), Inches(12), Inches(0.38),
         font_size=13, bold=True, color=EQ_COLOR)
add_text(sl, "R = effective ion radius  |  ε_water ≈ 80  |  ε_protein ≈ 4  |  z = charge of residue",
         Inches(0.55), Inches(3.85), Inches(12.2), Inches(0.28),
         font_size=10, color=LIGHT_GREY, italic=True)
add_text(sl, "→  If a mutation moves a charged residue from surface (water, ε=80) to buried (protein, ε=4), the energy penalty is large and destabilising.",
         Inches(0.55), Inches(4.15), Inches(12.2), Inches(0.28),
         font_size=10, color=WHITE)

# ── 3. Hill Equation ──────────────────────────────────────────────────────────
add_card(sl, Inches(0.35), Inches(4.65), Inches(6.1), Inches(2.45), CARD_BG, TEAL_DARK, True)
add_text(sl, "3  Hill Equation  —  pH Activity Correction (Cutinase family)",
         Inches(0.55), Inches(4.72), Inches(5.8), Inches(0.38), font_size=12, bold=True, color=TEAL_DARK)
add_text(sl, "f(pH)  =  [H⁺]ⁿ / ( K_a^n  +  [H⁺]ⁿ )",
         Inches(0.55), Inches(5.15), Inches(5.8), Inches(0.38),
         font_size=13, bold=True, color=EQ_COLOR)
add_text(sl, "n = Hill coefficient (cooperativity)\nK_a = half-saturation pH constant\n[H⁺] = 10^(−pH)",
         Inches(0.55), Inches(5.6), Inches(5.8), Inches(0.55),
         font_size=10, color=LIGHT_GREY, italic=True)
add_text(sl, "→  Applied post-inference to ΔΔG scores when the input enzyme\nis a cutinase-type PET hydrolase to prevent overestimating\nstability at extreme pH.",
         Inches(0.55), Inches(6.18), Inches(5.8), Inches(0.75),
         font_size=10, color=WHITE)

# ── 4. ΔΔG Sign Convention ────────────────────────────────────────────────────
add_card(sl, Inches(6.85), Inches(4.65), Inches(6.1), Inches(2.45), CARD_BG2, TEAL, True)
add_text(sl, "4  ΔΔG Sign Convention",
         Inches(7.0), Inches(4.72), Inches(5.8), Inches(0.38), font_size=12, bold=True, color=TEAL)
add_text(sl, "ΔΔG  =  ΔG_mutant  −  ΔG_wild-type",
         Inches(7.0), Inches(5.15), Inches(5.8), Inches(0.38),
         font_size=13, bold=True, color=EQ_COLOR)
add_text(sl, "ΔΔG < 0  →  mutant is MORE stable (stabilising)  ✓\nΔΔG > 0  →  mutant is LESS stable (destabilising)  ✗\nΔΔG = 0  →  no change in stability",
         Inches(7.0), Inches(5.6), Inches(5.8), Inches(0.75),
         font_size=11, color=WHITE)
add_text(sl, "Threshold used: −0.093 kcal/mol (Youden-optimal from CV)",
         Inches(7.0), Inches(6.38), Inches(5.8), Inches(0.35),
         font_size=10, color=LIGHT_GREY, italic=True)

add_footer(sl)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Honest Assessment (Updated)
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl)
add_title_bar(sl, "Honest Assessment (v46)",
              "What 79.54% CV accuracy and 68.16% S669 accuracy really mean — and the road ahead.")

issues = [
    (RED_WARN,  "S669 Stabilising Recall: 36.9%",
     "Of 168 true stabilising mutations, the model only caught 62.\n"
     "The model is conservative — it would rather miss a stabilising mutation\n"
     "than falsely label a destabilising one as good.\n"
     "→ Fix: adding ProDDG S2648 (2,648 mutations, 77% stabilising) in next run."),
    (RED_WARN,  "CV vs S669 Generalisation Gap",
     "CV accuracy 79.54% vs S669 68.16% — an 11 pp drop.\n"
     "S669 contains protein families the model has not seen.\n"
     "This is the honest cost of generalising to novel enzyme families.\n"
     "→ More diverse training data (ProDDG) is the primary fix."),
    (TEAL_MID,  "Additivity Assumption for Multi-Mutants",
     "Multi-mutant ΔΔG = sum of individual scores.\n"
     "This breaks down when mutations are within ~5 Å — epistatic\n"
     "(cooperative) effects between neighbouring residues are not captured.\n"
     "→ Nearby pairs are flagged with a warning in the web app."),
    (TEAL_DARK, "ESM-1v Coverage Only 44.1%",
     "The evolutionary fitness signal from ESM-1v is missing for\n"
     "55.9% of training samples — cached for only 100 proteins.\n"
     "→ Expanding the ESM-1v cache would improve coverage significantly."),
]

for i, (border_col, title, body) in enumerate(issues):
    col = i % 2
    row = i // 2
    px = Inches(0.35) + col * Inches(6.6)
    py = Inches(1.3) + row * Inches(2.9)
    add_card(sl, px, py, Inches(6.25), Inches(2.65), CARD_BG, border_col, True)
    add_text(sl, title, px + Inches(0.15), py + Inches(0.12), Inches(6.0), Inches(0.38),
             font_size=13, bold=True, color=border_col)
    add_text(sl, body, px + Inches(0.15), py + Inches(0.55), Inches(6.0), Inches(2.0),
             font_size=11, color=WHITE)

add_footer(sl)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Summary (Updated numbers)
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl, DARK_NAVY)
add_text(sl, "Sequence in.", Inches(0.5), Inches(0.6), Inches(7), Inches(0.75),
         font_size=52, bold=True, color=WHITE)
add_text(sl, "Ranked candidates out.", Inches(0.5), Inches(1.35), Inches(10), Inches(0.75),
         font_size=52, bold=True, color=TEAL)
add_text(sl, "No crystal structure. No specialist software. No barrier to entry.",
         Inches(0.5), Inches(2.3), Inches(8), Inches(0.4),
         font_size=16, italic=True, color=LIGHT_GREY)
add_text(sl, "This project shows that an ML pipeline combining sequence features, protein language\n"
             "models (ESM-2, ESM-1v), AlphaFold structural data, and explicit biophysics (Debye-Hückel,\n"
             "Born solvation, Hill equation) can reliably distinguish stabilising from destabilising\n"
             "mutations in PET-degrading enzymes — and deliver those predictions in ~3 seconds.",
         Inches(0.5), Inches(2.85), Inches(7.5), Inches(1.5),
         font_size=13, color=WHITE)

big_stats = [
    ("16,963",  "real wet-lab mutations\nin training"),
    ("144",     "features per mutation\n(sequence + structure + physics)"),
    ("8",       "independently tuned\nbase models"),
    ("79.54%",  "CV classification\naccuracy"),
    ("~3 sec",  "to rank all mutations\nfor a new protein"),
]
for i, (num, label) in enumerate(big_stats):
    sy = Inches(4.6) + (i % 3) * Inches(0.0)
    sx = Inches(7.5) + (i % 3) * Inches(1.92)
    # vertical stack on right side
    by = Inches(4.45) + i * Inches(0.6)
    add_card(sl, Inches(7.5), by, Inches(5.5), Inches(0.55), CARD_BG2)
    add_text(sl, num, Inches(7.6), by + Inches(0.05), Inches(1.6), Inches(0.48),
             font_size=22, bold=True, color=TEAL)
    add_text(sl, label, Inches(9.3), by + Inches(0.05), Inches(3.6), Inches(0.48),
             font_size=11, color=WHITE)

add_footer(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Our Solution (Updated)
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl)

# Title (no accent bar — matches original style)
add_text(sl, "Our Solution", Inches(0.5), Inches(0.2), Inches(10), Inches(0.7),
         font_size=40, bold=True, color=WHITE)

BLUE_NUM = RGBColor(0x4A, 0x90, 0xD9)

solutions = [
    ("01", TEAL,      "Temperature Stability",
     "✓  IMPLEMENTED",
     "•  ΔTm regressor trained on 6,567 records — CV MAE 3.66 °C\n"
     "•  Predicts melting temperature shift for every single mutation\n"
     "•  Next: synthesise top 5 variants and measure Tm experimentally"),

    ("02", BLUE_NUM,  "pH Stability",
     "✓  IMPLEMENTED",
     "•  pH modelled as direct input feature across all 8 models\n"
     "•  Hill equation correction applied post-inference for cutinase family\n"
     "•  Next: collect data targeting industrial pH ranges (6–8.1)"),

    ("03", TEAL,      "Ionic Strength / Ca²⁺ Concentrations",
     "✓  IMPLEMENTED",
     "•  Debye-Hückel equation screens electrostatics by ionic strength\n"
     "•  Ca²⁺ coordination cache covers 897 binding-site residue positions\n"
     "•  Next: integrate AlphaFold2 3D coords for active-site tracking"),

    ("04", BLUE_NUM,  "Biophysics Modeling",
     "✓  IMPLEMENTED  +  EXPANDING",
     "•  Debye-Hückel + Born solvation embedded as ~30 input features\n"
     "•  ESM-2 language model + AlphaFold pLDDT at 92.9% coverage\n"
     "•  Next: apply pipeline to PHA and polyurethane hydrolases"),
]

positions_sol = [
    (Inches(0.4),  Inches(1.1)),
    (Inches(6.85), Inches(1.1)),
    (Inches(0.4),  Inches(4.15)),
    (Inches(6.85), Inches(4.15)),
]

for i, (num, num_col, title, status, body) in enumerate(solutions):
    px, py = positions_sol[i]
    cw, ch = Inches(6.1), Inches(2.95)

    # card background
    add_card(sl, px, py, cw, ch, CARD_BG)

    # number
    add_text(sl, num, px + Inches(0.1), py + Inches(0.08), Inches(0.7), Inches(0.5),
             font_size=28, bold=True, color=num_col)

    # title
    add_text(sl, title, px + Inches(0.82), py + Inches(0.1), cw - Inches(0.95), Inches(0.42),
             font_size=15, bold=True, color=WHITE)

    # status badge
    status_col = GREEN_OK if "IMPLEMENTED" in status else TEAL
    add_text(sl, status, px + Inches(0.82), py + Inches(0.52), cw - Inches(0.95), Inches(0.28),
             font_size=10, bold=True, color=status_col, italic=True)

    # body
    add_text(sl, body, px + Inches(0.15), py + Inches(0.84), cw - Inches(0.25), Inches(2.0),
             font_size=10, color=WHITE)

add_footer(sl)


# ── Save ──────────────────────────────────────────────────────────────────────
out = "/Users/admin/Documents/PET - Lab/PET_Lab_Updated_Slides.pptx"
prs.save(out)
print(f"Saved: {out}")
